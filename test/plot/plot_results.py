"""
  macro for plotting analyzed jetscape events
  """

# This script plots histograms created in the analysis of Jetscape events
#
# Author: James Mulligan (james.mulligan@berkeley.edu)

# General
import os
import sys
import yaml
import argparse

# Data analysis and plotting
import ROOT
import ctypes
import numpy as np
import pptx             # pip install python-pptx

# Base class
sys.path.append('.')
from jetscape_analysis.base import common_base
from plot import plot_results_utils

# Prevent ROOT from stealing focus when plotting
ROOT.gROOT.SetBatch(True)

################################################################
class PlotResults(common_base.CommonBase):

    # ---------------------------------------------------------------
    # Constructor
    # ---------------------------------------------------------------
    def __init__(self, config_file='', input_file='', reference_file='', **kwargs):
        super(PlotResults, self).__init__(**kwargs)
        self.output_dir = os.path.dirname(input_file)
               
        self.plot_utils = plot_results_utils.PlotUtils()
        self.plot_utils.setOptions()
        ROOT.gROOT.ForceStyle()

        self.input_file = ROOT.TFile(input_file, 'READ')
        self.reference_file = ROOT.TFile(reference_file, 'READ')

        self.data_color = ROOT.kGray+3
        self.data_marker = 21
        self.jetscape_color = ROOT.kViolet-8
        self.jetscape_color_reference = ROOT.kTeal-8
        self.pre_post_ratio_color = ROOT.kRed-6
        self.jetscape_marker = 20
        self.alpha = 0.7
        self.marker_size = 1.5
        self.line_width = 2
        self.line_style = 1
        self.file_format = '.pdf'

        # Read config file
        with open(config_file, 'r') as stream:
            self.config = yaml.safe_load(stream)
        self.sqrts = self.config['sqrt_s']
        self.power = self.config['power']
        self.pt_ref = self.config['pt_ref']

        # We will write final results after all scalings, along with data, to file
        self.output_dict = {}
      
        print(self)

    #-------------------------------------------------------------------------------------------
    #-------------------------------------------------------------------------------------------
    #-------------------------------------------------------------------------------------------
    def plot_results(self):
    
        self.plot_hadron_observables(observable_type='hadron')
        
        self.plot_hadron_correlation_observables(observable_type='hadron_correlations')
        
        self.plot_jet_observables(observable_type='inclusive_chjet')
        
        if 'inclusive_jet' in self.config:
            self.plot_jet_observables(observable_type='inclusive_jet')
            
        if 'semi_inclusive_chjet' in self.config:
            self.plot_semi_inclusive_chjet_observables(observable_type='semi_inclusive_chjet')
            
        if 'dijet' in self.config:
            self.plot_jet_observables(observable_type='dijet')

        self.write_output_objects()
        
        # Generate pptx for convenience
        if self.file_format == '.png':
            self.generate_pptx()

    #-------------------------------------------------------------------------------------------
    # Plot hadron observables
    #-------------------------------------------------------------------------------------------
    def plot_hadron_observables(self, observable_type=''):
        print()
        print(f'Plot {observable_type} observables...')
                
        for observable, block in self.config[observable_type].items():
        
            if 'hepdata' not in block and 'custom_data' not in block:
                continue
        
            # Initialize observable configuration
            self.suffix = ''
            self.init_observable(observable_type, observable, block)
                
            # Plot observable
            self.plot_observable(observable_type, observable)

    #-------------------------------------------------------------------------------------------
    # Plot hadron correlation observables
    #-------------------------------------------------------------------------------------------
    def plot_hadron_correlation_observables(self, observable_type=''):
        print()
        print(f'Plot {observable_type} observables...')

        for observable, block in self.config[observable_type].items():

            if 'hepdata' not in block and 'custom_data' not in block:
                continue

            # Initialize observable configuration
            self.suffix = ''
            self.init_observable(observable_type, observable, block)
                
            # Histogram observable
            self.plot_observable(observable_type, observable)

    #-------------------------------------------------------------------------------------------
    # Histogram inclusive jet observables
    #-------------------------------------------------------------------------------------------
    def plot_jet_observables(self, observable_type=''):
        print()
        print(f'Plot {observable_type} observables...')

        for observable, block in self.config[observable_type].items():                
            for self.jet_R in block['jet_R']:
            
                # Optional: Loop through pt bins
                for pt_bin in range(len(block['pt'])-1):

                    if len(block['pt']) > 2:
                        pt_suffix = f'_pt{pt_bin}'
                    else:
                        pt_suffix = ''
                        
                    # Optional: subobservable
                    subobservable_label_list = ['']
                    if 'kappa' in block:
                        subobservable_label_list = [f'_k{kappa}' for kappa in block['kappa']]
                    for subobservable_label in subobservable_label_list:
                    
                        # Set normalization
                        self_normalize = False
                        for x in ['mass', 'g', 'ptd', 'charge', 'mg', 'zg', 'tg', 'xj']:
                            if x in observable:
                                self_normalize = True
                    
                        if 'SoftDrop' in block:
                            for grooming_setting in block['SoftDrop']:

                                print(f'      grooming_setting = {grooming_setting}')
                                zcut = grooming_setting['zcut']
                                beta = grooming_setting['beta']
                                
                                self.suffix = f'_R{self.jet_R}_zcut{zcut}_beta{beta}{subobservable_label}'
                                if 'hepdata' not in block and 'custom_data' not in block:
                                    continue
                    
                                # Initialize observable configuration
                                self.init_observable(observable_type, observable, block, pt_suffix=pt_suffix, self_normalize=self_normalize)
                            
                                # Plot observable
                                self.plot_observable(observable_type, observable, pt_suffix)
                            
                        else:

                            self.suffix = f'_R{self.jet_R}{subobservable_label}'
                            if 'hepdata' not in block and 'custom_data' not in block:
                                continue

                            # Initialize observable configuration
                            self.init_observable(observable_type, observable, block, pt_suffix=pt_suffix, self_normalize=self_normalize)
                        
                            # Plot observable
                            self.plot_observable(observable_type, observable, pt_suffix)

    #-------------------------------------------------------------------------------------------
    # Histogram semi-inclusive jet observables
    #-------------------------------------------------------------------------------------------
    def plot_semi_inclusive_chjet_observables(self, observable_type=''):
        print()
        print(f'Plot {observable_type} observables...')
        
        for observable, block in self.config[observable_type].items():                    
            for self.jet_R in block['jet_R']:
                
                self.suffix = f'_R{self.jet_R}'
                if 'hepdata' not in block and 'custom_data' not in block:
                    continue
                    
                # Set normalization
                self_normalize = False
                for x in ['nsubjettiness']:
                    if x in observable:
                        self_normalize = True

                # Initialize observable configuration
                self.init_observable(observable_type, observable, block, self_normalize=self_normalize)
            
                # Plot observable
                self.plot_observable(observable_type, observable)

    #-------------------------------------------------------------------------------------------
    # Initialize a single observable's config
    #-------------------------------------------------------------------------------------------
    def init_observable(self, observable_type, observable, block, pt_suffix='', self_normalize=False):
    
        # Initialize an empty dict containing relevant info
        self.observable_settings = {}

        # Initialize common settings into class members
        self.init_common_settings(observable, block)
        
        #-----------------------------------------------------------
        # Initialize data distribution into self.observable_settings
        if 'hepdata' in block:
            self.observable_settings['data_distribution'] = self.plot_utils.tgraph_from_hepdata(block, self.sqrts, observable_type, observable, suffix=self.suffix, pt_suffix=pt_suffix)
        elif 'custom_data' in block:
            self.observable_settings['data_distribution'] = self.plot_utils.tgraph_from_yaml(block, self.sqrts, observable_type, observable, suffix=self.suffix, pt_suffix=pt_suffix)
        else:
            self.observable_settings['data_distribution'] = None

        #-----------------------------------------------------------
        # Initialize JETSCAPE distribution into self.observable_settings
        self.initialize_jetscape_distribution(observable_type, observable, pt_suffix=pt_suffix, self_normalize=self_normalize, reference=False) # pre
        self.initialize_jetscape_distribution(observable_type, observable, pt_suffix=pt_suffix, self_normalize=self_normalize, reference=True)  # post

        #-----------------------------------------------------------
        # Form ratio of JETSCAPE to data and load into self.observable_settings
        if self.observable_settings['data_distribution'] and self.observable_settings['jetscape_distribution'] and not self.observable_settings['jetscape_distribution'].InheritsFrom(ROOT.TH2.Class()) and not self.skip_pp_ratio:
            self.observable_settings['ratio'] = self.plot_utils.divide_histogram_by_tgraph(self.observable_settings['jetscape_distribution'],
                                                                                           self.observable_settings['data_distribution'])
            self.observable_settings['ratio_reference'] = self.plot_utils.divide_histogram_by_tgraph(self.observable_settings['jetscape_distribution_reference'],
                                                                                                     self.observable_settings['data_distribution'])
        else:
            self.observable_settings['ratio'] = None
            self.observable_settings['ratio_reference'] = None

    #-------------------------------------------------------------------------------------------
    # Initialize from settings from config file into class members
    #-------------------------------------------------------------------------------------------
    def init_common_settings(self, observable, block):

        self.xtitle = block['xtitle']
        if 'eta_cut' in block:
            self.eta_cut = block['eta_cut']
        if 'y_cut' in block:
            self.y_cut = block['y_cut']
        if 'pt' in block:
            self.pt = block['pt']
        if 'eta_cut_R' in block:
            self.eta_R = block['eta_cut_R']
            self.eta_cut = np.round(self.eta_R - self.jet_R, decimals=1)
        if 'c_ref' in block:
            index = block['jet_R'].index(self.jet_R)
            self.c_ref = block['c_ref'][index]
        if 'low_trigger_range' in block:
            self.low_trigger_range = block['low_trigger_range']
        if 'high_trigger_range' in block:
            self.high_trigger_range = block['high_trigger_range']
        if 'trigger_range' in block:
            self.trigger_range = block['trigger_range']
        if 'logy' in block:
            self.logy = block['logy']
        else:
            self.logy = False
                 
        if 'ytitle_pp' in block:
            self.ytitle = block['ytitle_pp']
        else:
            self.ytitle = ''
        if 'y_min_pp' in block:
            self.y_min = float(block['y_min_pp'])
            self.y_max = float(block['y_max_pp'])
        else:
            self.y_min = 0.
            self.y_max = 1.99
        if 'y_ratio_min' in block:
            self.y_ratio_min = block['y_ratio_min']
            self.y_ratio_max = block['y_ratio_max']
        else:
            self.y_ratio_min = 0.
            self.y_ratio_max = 1.99

        if 'skip_pp' in block:
            self.skip_pp = block['skip_pp']
        else:
            self.skip_pp = False
        if 'skip_pp_ratio' in block:
            self.skip_pp_ratio = block['skip_pp_ratio']
        else:
            self.skip_pp_ratio = False
        if 'scale_by' in block:
            self.scale_by = block['scale_by']
        else:
            self.scale_by = None

    #-------------------------------------------------------------------------------------------
    # Initialize JETSCAPE distribution into self.observable_settings
    #-------------------------------------------------------------------------------------------
    def initialize_jetscape_distribution(self, observable_type, observable, pt_suffix='', self_normalize=False, reference=False):

        if reference:
            jetscape_key = 'jetscape_distribution_reference'
        else:
            jetscape_key = 'jetscape_distribution'

        # Get histogram and add to self.observable_settings
        #  - In AA case, also add hole histogram
        #  - In the case of semi-inclusive measurements construct difference of histograms
        self.get_histogram(observable_type, observable, pt_suffix=pt_suffix, jetscape_key=jetscape_key)

        # Normalization
        # Note: If we divide by the sum of weights (corresponding to n_events) and multiply by the
        #       pt-hat cross-section, then JETSCAPE distribution gives cross-section: dsigma/dx (in mb)
        self.scale_histogram(self.observable_settings[jetscape_key], observable_type, observable, 
                             pt_suffix=pt_suffix, self_normalize=self_normalize)
    
    #-------------------------------------------------------------------------------------------
    # Get histogram and add to self.observable_settings
    #  - In AA case, also add hole histogram
    #  - In the case of semi-inclusive measurements construct difference of histograms    
    #-------------------------------------------------------------------------------------------
    def get_histogram(self, observable_type, observable, pt_suffix='', jetscape_key=''):

        keys = [key.ReadObj().GetTitle() for key in self.input_file.GetListOfKeys()]

        # In the case of semi-inclusive measurements construct difference of histograms  
        if 'semi_inclusive' in observable_type:
            self.construct_semi_inclusive_histogram(keys, observable_type, observable, jetscape_key)

        # For all other histograms, get the histogram directly
        # For those observables that we do hole subtraction, also construct the hole histogram
        else:

            # Get histogram
            self.hname = f'h_{observable_type}_{observable}{self.suffix}{pt_suffix}'
            if self.hname in keys:
                h_jetscape = self.input_file.Get(self.hname)
                h_jetscape.SetDirectory(0)
                h_jetscape.SetName(f'{self.hname}_{jetscape_key}')
            else:
                h_jetscape = None
            self.observable_settings[jetscape_key] = h_jetscape

    #-------------------------------------------------------------------------------------------
    # Construct semi-inclusive observables from difference of histograms
    #-------------------------------------------------------------------------------------------
    def construct_semi_inclusive_histogram(self, keys, observable_type, observable, jetscape_key):

        if self.sqrts == 2760: # Delta recoil

            hname_low_trigger = f'h_{observable_type}_{observable}_R{self.jet_R}_lowTrigger'
            hname_high_trigger = f'h_{observable_type}_{observable}_R{self.jet_R}_highTrigger'
            hname_ntrigger = f'h_{observable_type}_alice_trigger_pt{observable}'
            self.hname = f'h_{observable_type}_{observable}_R{self.jet_R}'
            if hname_low_trigger in keys and hname_high_trigger in keys and hname_ntrigger in keys:
                h_jetscape_low = self.input_file.Get(hname_low_trigger)
                h_jetscape_low.SetDirectory(0)
                h_jetscape_high = self.input_file.Get(hname_high_trigger)
                h_jetscape_high.SetDirectory(0)
                h_jetscape_ntrigger = self.input_file.Get(hname_ntrigger)
                h_jetscape_ntrigger.SetDirectory(0)
                
                low_trigger = (self.low_trigger_range[0]+self.low_trigger_range[1])/2
                high_trigger = (self.high_trigger_range[0]+self.high_trigger_range[1])/2
                n_trig_high = h_jetscape_ntrigger.GetBinContent(h_jetscape_ntrigger.FindBin(high_trigger))
                n_trig_low = h_jetscape_ntrigger.GetBinContent(h_jetscape_ntrigger.FindBin(low_trigger))
                h_jetscape_high.Scale(1./n_trig_high)
                h_jetscape_low.Scale(1./n_trig_low) 
                self.observable_settings[jetscape_key] = h_jetscape_high.Clone('h_delta_recoil')
                self.observable_settings[jetscape_key].Add(h_jetscape_low, -1)
            else:
                self.observable_settings[jetscape_key] = None

        elif self.sqrts == 200:

            hname = f'h_{observable_type}_{observable}_R{self.jet_R}'
            hname_ntrigger = f'h_{observable_type}_star_trigger_pt'
            self.hname = f'h_{observable_type}_{observable}_R{self.jet_R}'
            if hname in keys and hname_ntrigger in keys:
                self.observable_settings[jetscape_key] = self.input_file.Get(hname)
                self.observable_settings[jetscape_key].SetDirectory(0)
                h_jetscape_ntrigger = self.input_file.Get(hname_ntrigger)
                h_jetscape_ntrigger.SetDirectory(0)
                
                trigger = (self.trigger_range[0]+self.trigger_range[1])/2
                n_trig = h_jetscape_ntrigger.GetBinContent(h_jetscape_ntrigger.FindBin(trigger))
                self.observable_settings[jetscape_key].Scale(1./n_trig)
            else:
                self.observable_settings[jetscape_key] = None

    #-------------------------------------------------------------------------------------------
    # Scale histogram for:
    #  - xsec/weight_sum
    #  - bin width
    #  - observable-specific factors (eta, sigma_inel, n_jets, ...)
    #  - self-normalization
    #-------------------------------------------------------------------------------------------
    def scale_histogram(self, h, observable_type, observable, pt_suffix='', self_normalize=False):

        if not h:
            return

        #--------------------------------------------------
        # (1) Scale all histograms by the min-pt-hat cross-section and weight-sum
        h_xsec = self.input_file.Get('h_xsec')
        weight_sum = self.input_file.Get('h_weight_sum').GetBinContent(1)
        xsec = h_xsec.GetBinContent(1) / h_xsec.GetEntries()
        h.Scale(xsec/weight_sum)
        
        #--------------------------------------------------
        # (2) Scale all histograms by bin width
        h.Scale(1., 'width')

        #--------------------------------------------------
        # (3) Scale for observable-specific factors
        if self.sqrts == 200:
            if observable_type == 'hadron':
                if observable == 'pt_ch_star':
                    sigma_inel = 42.
                    h.Scale(1./(2*self.eta_cut))
                    h.Scale(1./(2*np.pi))
                    h.Scale(1./sigma_inel)
                if observable == 'pt_pi0_phenix':
                    sigma_inel = 42. # Update
                    h.Scale(1./(2*self.eta_cut))
                    h.Scale(1./(2*np.pi))
                    h.Scale(1./sigma_inel)
            elif observable_type == 'inclusive_chjet':
                if observable == 'pt_star':
                    h.Scale(1./(2*self.eta_cut))
            elif observable_type == 'semi_inclusive_chjet':
                # Note that n_trig histograms should also be scaled by xsec/n_events
                if observable in ['IAA_star', 'dphi_star']:
                    h.Scale(1./(xsec/weight_sum))
        
        if self.sqrts == 2760:
            if observable_type == 'hadron':
                if observable in ['pt_ch_alice', 'pt_pi_alice', 'pt_pi0_alice']:
                    sigma_inel = 61.8
                    h.Scale(1./(2*self.eta_cut))
                    h.Scale(1./sigma_inel)
                elif observable == 'pt_ch_atlas':
                    h.Scale(1./(2*self.eta_cut))
                    h.Scale(1./(2*np.pi))
                elif observable == 'pt_ch_cms':
                    sigma_inel = 64. 
                    h.Scale(1./(2*self.eta_cut))
                    h.Scale(1./(2*np.pi))
                    h.Scale(1./sigma_inel)
            elif observable_type == 'inclusive_jet':
                if observable == 'pt_alice':
                    sigma_inel = 62.1
                    h.Scale(1./(2*self.eta_cut))
                    h.Scale(1./sigma_inel)
                if observable == 'pt_atlas':
                    h.Scale(1./(2*self.y_cut))
                    h.Scale(1.e6) # convert to nb
                if observable == 'pt_cms':
                    h.Scale(1./(2*self.eta_cut))
                    h.Scale(1.e6) # convert to nb
                if observable in ['Dz_atlas', 'Dpt_atlas', 'Dz_cms', 'Dpt_cms']:
                    if observable in ['Dz_atlas', 'Dpt_atlas']:
                        hname = f'h_{observable_type}_Dz_atlas{self.suffix}_Njets{pt_suffix}'
                    elif observable in ['Dz_cms', 'Dpt_cms']:
                        hname = f'h_{observable_type}_Dz_cms{self.suffix}_Njets{pt_suffix}'
                    h_njets = self.input_file.Get(hname)
                    h_njets.SetDirectory(0)
                    n_jets = h_njets.GetBinContent(1) # Note that Njets histogram should also be scaled by xsec/n_events
                    if n_jets > 0.:
                        h.Scale(1./(n_jets * (xsec/weight_sum)))
                    else:
                        print('WARNING: N_jets = 0')
            elif observable_type == 'inclusive_chjet':
                if observable == 'pt_alice':
                    h.Scale(1./(2*self.eta_cut))
            elif observable_type == 'semi_inclusive_chjet':
                # Note that n_trig histograms should also be scaled by xsec/n_events
                if observable in ['IAA_alice', 'dphi_alice']:
                    h.Scale(1./(xsec/weight_sum))
                        
        if self.sqrts == 5020:
            if observable_type == 'hadron':
                if observable in ['pt_ch_alice', 'pt_pi_alice']:
                    sigma_inel = 67.6
                    h.Scale(1./(2*self.eta_cut))
                    h.Scale(1./sigma_inel)
                elif observable == 'pt_ch_cms':
                    sigma = 70.
                    h.Scale(1./(2*self.eta_cut))
                    h.Scale(1./(2*np.pi))
                    h.Scale(1./sigma)
            elif observable_type == 'inclusive_jet':
                if observable == 'pt_alice':
                    h.Scale(1./(2*self.eta_cut))
                if observable == 'pt_atlas':
                    h.Scale(1./(2*self.y_cut))
                    h.Scale(1.e6) # convert to nb
                if observable == 'pt_cms':
                    h.Scale(1./(2*self.eta_cut))
                    h.Scale(1.e6) # convert to nb
                if observable in ['Dz_atlas', 'Dpt_atlas']:
                    hname = f'h_{observable_type}_Dz_atlas{self.suffix}_Njets{pt_suffix}'
                    h_njets = self.input_file.Get(hname)
                    h_njets.SetDirectory(0)
                    n_jets = h_njets.GetBinContent(1) # Note that Njets histogram should also be scaled by xsec/n_events
                    if n_jets > 0.:
                        h.Scale(1./(n_jets * (xsec/weight_sum)))
                    else:
                        print('WARNING: N_jets = 0')

        # Scale by bin-dependent factor (approximation for pp QA)
        if self.scale_by:
            nBins = h.GetNbinsX()
            for bin in range(1, nBins+1):
                h_x = h.GetBinCenter(bin)
                h_y = h.GetBinContent(bin)
                h_error = h.GetBinError(bin)
                if self.scale_by == '1/pt':
                    scale_factor = 1./h_x
                h.SetBinContent(bin, h_y*scale_factor)
                h.SetBinError(bin, h_error*scale_factor)

        #--------------------------------------------------
        # (4) Self-normalize histogram if appropriate
        if self_normalize:
            if observable in ['zg_alice', 'tg_alice']:
                min_bin = 0
            else:
                min_bin = 1
            nbins = h.GetNbinsX()
            integral = h.Integral(min_bin, nbins, 'width')
            if integral > 0.:
                h.Scale(1./integral)
            else:
                print('WARNING: integral for self-normalization is 0')

    #-------------------------------------------------------------------------------------------
    # Plot distributions in upper panel, and ratio in lower panel
    #-------------------------------------------------------------------------------------------
    def plot_observable(self, observable_type, observable, pt_suffix='', logy = False):
    
        if not self.observable_settings['jetscape_distribution']:
            return
        
        label = f'{observable_type}_{observable}_{self.sqrts}_{self.suffix}_{pt_suffix}'

        # Plot distribution, and ratio to data
        self.plot_distribution_and_ratio(observable_type, observable, label, pt_suffix=pt_suffix)

    #-------------------------------------------------------------------------------------------
    # Plot distributions in upper panel, and ratio in lower panel
    #-------------------------------------------------------------------------------------------
    def plot_distribution_and_ratio(self, observable_type, observable, label, pt_suffix=''):
            
        c = ROOT.TCanvas('c', 'c', 600, 650)
        c.Draw()
        c.cd()
        
        # Distribution
        pad2_dy = 0.45
        pad1 = ROOT.TPad('myPad', 'The pad',0,pad2_dy,1,1)
        pad1.SetLeftMargin(0.2)
        pad1.SetTopMargin(0.08)
        pad1.SetRightMargin(0.04)
        pad1.SetBottomMargin(0.)
        pad1.SetTicks(0,1)
        pad1.Draw()
        if self.logy:
            pad1.SetLogy()
        pad1.cd()
        
        legend = ROOT.TLegend(0.58,0.55,0.75,0.73)
        self.plot_utils.setup_legend(legend, 0.055, sep=-0.1)

        legend_ratio = ROOT.TLegend(0.45,0.85,0.65,0.98)
        self.plot_utils.setup_legend(legend_ratio, 0.06, sep=-0.1)
            
        self.bins = np.array(self.observable_settings['jetscape_distribution'].GetXaxis().GetXbins())
        myBlankHisto = ROOT.TH1F('myBlankHisto','Blank Histogram', 1, self.bins[0], self.bins[-1])
        myBlankHisto.SetNdivisions(505)
        myBlankHisto.SetXTitle(self.xtitle)
        myBlankHisto.SetYTitle(self.ytitle)
        myBlankHisto.SetMaximum(self.y_max)
        myBlankHisto.SetMinimum(self.y_min) # Don't draw 0 on top panel
        myBlankHisto.GetYaxis().SetTitleSize(0.08)
        myBlankHisto.GetYaxis().SetTitleOffset(1.1)
        myBlankHisto.GetYaxis().SetLabelSize(0.06)
        myBlankHisto.Draw('E')

        # Ratio
        c.cd()
        pad2 = ROOT.TPad('pad2', 'pad2', 0, 0.02, 1, pad2_dy)
        pad2.SetTopMargin(0)
        pad2.SetBottomMargin(0.21)
        pad2.SetLeftMargin(0.2)
        pad2.SetRightMargin(0.04)
        pad2.SetTicks(0,1)
        pad2.Draw()
        pad2.cd()
              
        myBlankHisto2 = myBlankHisto.Clone('myBlankHisto_C')
        myBlankHisto2.SetYTitle('Ratio')
        myBlankHisto2.SetXTitle(self.xtitle)
        myBlankHisto2.GetXaxis().SetTitleSize(26)
        myBlankHisto2.GetXaxis().SetTitleFont(43)
        myBlankHisto2.GetXaxis().SetTitleOffset(2.3)
        myBlankHisto2.GetXaxis().SetLabelFont(43)
        myBlankHisto2.GetXaxis().SetLabelSize(22)
        myBlankHisto2.GetYaxis().SetTitleSize(28)
        myBlankHisto2.GetYaxis().SetTitleFont(43)
        myBlankHisto2.GetYaxis().SetTitleOffset(2.)
        myBlankHisto2.GetYaxis().SetLabelFont(43)
        myBlankHisto2.GetYaxis().SetLabelSize(20)
        myBlankHisto2.GetYaxis().SetNdivisions(505)
        myBlankHisto2.GetYaxis().SetRangeUser(self.y_ratio_min, self.y_ratio_max)
        myBlankHisto2.Draw('')
                
        # Draw data
        pad1.cd()
        if self.observable_settings['data_distribution']:
            self.output_dict[f'data_distribution_{label}'] = self.observable_settings['data_distribution']
            self.observable_settings['data_distribution'].SetName(f'data_distribution_{label}')
            self.observable_settings['data_distribution'].SetMarkerSize(self.marker_size)
            self.observable_settings['data_distribution'].SetMarkerStyle(self.data_marker)
            self.observable_settings['data_distribution'].SetMarkerColor(self.data_color)
            self.observable_settings['data_distribution'].SetLineStyle(self.line_style)
            self.observable_settings['data_distribution'].SetLineWidth(self.line_width)
            self.observable_settings['data_distribution'].SetLineColor(self.data_color)
            self.observable_settings['data_distribution'].Draw('PE Z same')
            legend.AddEntry(self.observable_settings['data_distribution'], 'Data', 'PE')
        
        # Draw JETSCAPE
        self.output_dict[f'jetscape_distribution_{label}_post'] = self.observable_settings['jetscape_distribution']
        if self.observable_settings['jetscape_distribution'].GetNbinsX() > 1:
            self.observable_settings['jetscape_distribution'].SetFillColor(self.jetscape_color)
            self.observable_settings['jetscape_distribution'].SetFillColorAlpha(self.jetscape_color, self.alpha)
            self.observable_settings['jetscape_distribution'].SetFillStyle(1001)
            self.observable_settings['jetscape_distribution'].SetMarkerSize(0.)
            self.observable_settings['jetscape_distribution'].SetMarkerStyle(0)
            self.observable_settings['jetscape_distribution'].SetLineWidth(0)
            self.observable_settings['jetscape_distribution'].DrawCopy('E3 same')
        elif self.observable_settings['jetscape_distribution'].GetNbinsX() == 1:
            self.observable_settings['jetscape_distribution'].SetMarkerSize(self.marker_size)
            self.observable_settings['jetscape_distribution'].SetMarkerStyle(self.data_marker+1)
            self.observable_settings['jetscape_distribution'].SetMarkerColor(self.jetscape_color)
            self.observable_settings['jetscape_distribution'].SetLineStyle(self.line_style)
            self.observable_settings['jetscape_distribution'].SetLineWidth(self.line_width)
            self.observable_settings['jetscape_distribution'].SetLineColor(self.jetscape_color)
            self.observable_settings['jetscape_distribution'].DrawCopy('PE same')
        legend.AddEntry(self.observable_settings['jetscape_distribution'], 'JETSCAPE (post)', 'f')

        # Draw JETSCAPE reference
        self.output_dict[f'jetscape_distribution_{label}_pre'] = self.observable_settings['jetscape_distribution_reference']
        if self.observable_settings['jetscape_distribution_reference'].GetNbinsX() > 1:
            self.observable_settings['jetscape_distribution_reference'].SetFillColor(self.jetscape_color_reference)
            self.observable_settings['jetscape_distribution_reference'].SetFillColorAlpha(self.jetscape_color_reference, self.alpha)
            self.observable_settings['jetscape_distribution_reference'].SetFillStyle(3144)
            self.observable_settings['jetscape_distribution_reference'].SetMarkerSize(0.)
            self.observable_settings['jetscape_distribution_reference'].SetMarkerStyle(0)
            self.observable_settings['jetscape_distribution_reference'].SetLineWidth(0)
            self.observable_settings['jetscape_distribution_reference'].DrawCopy('E3 same')
        elif self.observable_settings['jetscape_distribution_reference'].GetNbinsX() == 1:
            self.observable_settings['jetscape_distribution_reference'].SetMarkerSize(self.marker_size)
            self.observable_settings['jetscape_distribution_reference'].SetMarkerStyle(self.data_marker+1)
            self.observable_settings['jetscape_distribution_reference'].SetMarkerColor(self.jetscape_color_reference)
            self.observable_settings['jetscape_distribution_reference'].SetLineStyle(self.line_style)
            self.observable_settings['jetscape_distribution_reference'].SetLineWidth(self.line_width)
            self.observable_settings['jetscape_distribution_reference'].SetLineColor(self.jetscape_color_reference)
            self.observable_settings['jetscape_distribution_reference'].DrawCopy('PE same')
        legend.AddEntry(self.observable_settings['jetscape_distribution_reference'], 'JETSCAPE (pre)', 'f')

        legend.Draw()
        
        # Draw ratio

        # Data
        pad2.cd()
        draw_data_ratio = False
        if draw_data_ratio:
            if self.observable_settings['data_distribution']:
                data_ratio = self.plot_utils.divide_tgraph_by_tgraph(self.observable_settings['data_distribution'],
                                                                    self.observable_settings['data_distribution'])
                data_ratio.Draw('PE Z same')
                self.output_dict[f'data_ratio_{label}'] = data_ratio

        draw_jetscape_data_ratio = False
        if draw_jetscape_data_ratio:

            # JETSCAPE
            if self.observable_settings['ratio']:
                self.output_dict[f'ratio_{label}'] = self.observable_settings['ratio']
                if self.observable_settings['ratio'].GetN() > 1:
                    self.observable_settings['ratio'].SetFillColor(self.jetscape_color)
                    self.observable_settings['ratio'].SetFillColorAlpha(self.jetscape_color, self.alpha)
                    self.observable_settings['ratio'].SetFillStyle(1001)
                    self.observable_settings['ratio'].SetMarkerSize(0.)
                    self.observable_settings['ratio'].SetMarkerStyle(0)
                    self.observable_settings['ratio'].SetLineWidth(0)
                    self.observable_settings['ratio'].Draw('E3 same')
                elif self.observable_settings['ratio'].GetN() == 1:
                    self.observable_settings['ratio'].SetMarkerSize(self.marker_size)
                    self.observable_settings['ratio'].SetMarkerStyle(self.data_marker+1)
                    self.observable_settings['ratio'].SetMarkerColor(self.jetscape_color)
                    self.observable_settings['ratio'].SetLineStyle(self.line_style)
                    self.observable_settings['ratio'].SetLineWidth(self.line_width)
                    self.observable_settings['ratio'].SetLineColor(self.jetscape_color)
                    self.observable_settings['ratio'].Draw('PE same')

            # JETSCAPE reference
            if self.observable_settings['ratio_reference']:
                self.output_dict[f'ratio_reference_{label}'] = self.observable_settings['ratio_reference']
                if self.observable_settings['ratio_reference'].GetN() > 1:
                    self.observable_settings['ratio_reference'].SetFillColor(self.jetscape_color_reference)
                    self.observable_settings['ratio_reference'].SetFillColorAlpha(self.jetscape_color_reference, self.alpha)
                    self.observable_settings['ratio_reference'].SetFillStyle(3144)
                    self.observable_settings['ratio_reference'].SetMarkerSize(0.)
                    self.observable_settings['ratio_reference'].SetMarkerStyle(0)
                    self.observable_settings['ratio_reference'].SetLineWidth(0)
                    self.observable_settings['ratio_reference'].Draw('E3 same')
                elif self.observable_settings['ratio_reference'].GetN() == 1:
                    self.observable_settings['ratio_reference'].SetMarkerSize(self.marker_size)
                    self.observable_settings['ratio_reference'].SetMarkerStyle(self.data_marker+1)
                    self.observable_settings['ratio_reference'].SetMarkerColor(self.jetscape_color_reference)
                    self.observable_settings['ratio_reference'].SetLineStyle(self.line_style)
                    self.observable_settings['ratio_reference'].SetLineWidth(self.line_width)
                    self.observable_settings['ratio_reference'].SetLineColor(self.jetscape_color_reference)
                    self.observable_settings['ratio_reference'].Draw('PE same')

        # JETSCAPE post/pre
        if self.observable_settings['jetscape_distribution']:
            h_ratio = self.observable_settings['jetscape_distribution'].Clone('pre_post_ratio')
            h_ratio.Divide(self.observable_settings['jetscape_distribution_reference'])
            if h_ratio.GetNbinsX() > 1:
                h_ratio.SetFillColor(self.pre_post_ratio_color)
                h_ratio.SetFillColorAlpha(self.pre_post_ratio_color, self.alpha)
                h_ratio.SetFillStyle(1001)
                h_ratio.SetMarkerSize(0.)
                h_ratio.SetMarkerStyle(0)
                h_ratio.SetLineWidth(0)
                h_ratio.Draw('E3 same')
            elif h_ratio.GetNbinsX() == 1:
                h_ratio.SetMarkerSize(self.marker_size)
                h_ratio.SetMarkerStyle(self.data_marker+1)
                h_ratio.SetMarkerColor(self.pre_post_ratio_color)
                h_ratio.SetLineStyle(self.line_style)
                h_ratio.SetLineWidth(self.line_width)
                h_ratio.SetLineColor(self.pre_post_ratio_color)
                h_ratio.Draw('PE same')

        if self.observable_settings['ratio']:
            if draw_jetscape_data_ratio:
                legend_ratio.AddEntry(self.observable_settings['ratio'], 'JETSCAPE/Data (post)', 'f')
                legend_ratio.AddEntry(self.observable_settings['ratio_reference'], 'JETSCAPE/Data (pre)', 'f')
        if draw_data_ratio:
            if self.observable_settings['data_distribution']:
                legend_ratio.AddEntry(data_ratio, 'Data uncertainties', 'PE')
        
        legend_ratio.AddEntry(h_ratio, 'JETSCAPE post/pre', 'f')
        legend_ratio.Draw()
        
        line = ROOT.TLine(self.bins[0], 1, self.bins[-1], 1)
        line.SetLineColor(920+2)
        line.SetLineStyle(2)
        line.SetLineWidth(2)
        line.Draw()
        
        pad1.cd()
        text_latex = ROOT.TLatex()
        text_latex.SetNDC()
        
        x = 0.25
        text_latex.SetTextSize(0.065)
        text = f'#bf{{{observable_type}_{observable}}} #sqrt{{#it{{s}}}} = {self.sqrts/1000.} TeV'
        text_latex.DrawLatex(x, 0.83, text)
        text = f'{self.suffix} {pt_suffix}'
        text_latex.DrawLatex(x, 0.73, text)

        if self.skip_pp and not self.observable_settings['data_distribution']:
            text = 'skip data plot -- no pp data in HEPData'
            text_latex.DrawLatex(x, 0.53, text)

        pad2.cd()
        if draw_jetscape_data_ratio and self.skip_pp_ratio:
            text = 'skip ratio plot -- pp/AA binning mismatch'
            text_latex.DrawLatex(x, 0.93, text)

        if self.observable_settings['data_distribution']:
            c.SaveAs(os.path.join(self.output_dir, f'{self.hname}{self.file_format}'))
        c.Close()

    # ---------------------------------------------------------------
    # Save all ROOT histograms to file
    # ---------------------------------------------------------------
    def write_output_objects(self):

        # Save output objects
        self.output_filename = os.path.join(self.output_dir, 'final_results.root')
        fout = ROOT.TFile(self.output_filename, 'recreate')
        fout.cd()
        for key,val in self.output_dict.items():
            val.SetName(key)
            val.Write()
            if isinstance(val, (ROOT.TH1)):
                val.SetDirectory(0)
                del val
        fout.Close()

    #-------------------------------------------------------------------------------------------
    # Generate pptx of one plot per slide, for convenience
    #-------------------------------------------------------------------------------------------
    def generate_pptx(self):
    
        # Create a blank presentation
        p = pptx.Presentation()
        
        # Set slide layouts
        title_slide_layout = p.slide_layouts[0]
        blank_slide_layout = p.slide_layouts[6]
        
        # Make a title slide
        slide = p.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        title.text = f'QA for {self.sqrts/1000.} TeV'
        author = slide.placeholders[1]
        author.text = 'STAT WG'
        
        # Loop through all output files and plot
        files = [f for f in os.listdir(self.output_dir) if f.endswith(self.file_format)]
        for file in sorted(files):
            img = os.path.join(self.output_dir, file)
            slide = p.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(img, left=pptx.util.Inches(2.),
                                          top=pptx.util.Inches(1.),
                                          width=pptx.util.Inches(5.))
            
        p.save(os.path.join(self.output_dir, 'results.pptx'))

#-------------------------------------------------------------------------------------------
#-------------------------------------------------------------------------------------------
#-------------------------------------------------------------------------------------------
if __name__ == '__main__':
    print('Executing plot_results_STAT.py...')
    print('')

    # Define arguments
    parser = argparse.ArgumentParser(description='Plot JETSCAPE events')
    parser.add_argument(
        '-c',
        '--configFile',
        action='store',
        type=str,
        metavar='configFile',
        default='config/TG3.yaml',
        help='Config file'
    )
    parser.add_argument(
        '-i',
        '--inputFile',
        action='store',
        type=str,
        metavar='inputFile',
        default='pp_5020_plot/histograms_5020_merged.root',
        help='Input file'
    )
    parser.add_argument(
        '-r',
        '--refFile',
        action='store',
        type=str,
        metavar='refFile',
        default='final_results.root',
        help='pp reference file'
    )

    # Parse the arguments
    args = parser.parse_args()
    
    # If invalid configFile is given, exit
    if not os.path.exists(args.configFile):
        print('File "{0}" does not exist! Exiting!'.format(args.configFile))
        sys.exit(0)

    # If invalid inputDir is given, exit
    if not os.path.exists(args.inputFile):
        print('File "{0}" does not exist! Exiting!'.format(args.inputFile))
        sys.exit(0)

    analysis = PlotResults(config_file=args.configFile, input_file=args.inputFile, reference_file=args.refFile)
    analysis.plot_results()
